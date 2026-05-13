"""PPTX export — python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def build_pptx(path: Path, title: str, sections: list[dict]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.bold = True

    for sec in sections:
        s = prs.slides.add_slide(blank)
        # Heading
        head = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        head.text_frame.text = sec.get("heading", "")
        head.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        head.text_frame.paragraphs[0].runs[0].font.bold = True

        # Body
        if sec.get("markdown"):
            body = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
            tf = body.text_frame
            tf.word_wrap = True
            for i, line in enumerate(sec["markdown"].splitlines()):
                if i == 0:
                    para = tf.paragraphs[0]
                else:
                    para = tf.add_paragraph()
                para.text = line
                if para.runs:
                    para.runs[0].font.size = Pt(16)

    prs.save(path)
